from __future__ import annotations

from pathlib import Path
from math import cos, sin, pi
import random

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/kane/Documents/agent-passport")
OUT_DIR = ROOT / "docs" / "generated"
ASSET_DIR = OUT_DIR / "assets"
OUT_DIR.mkdir(parents=True, exist_ok=True)
ASSET_DIR.mkdir(parents=True, exist_ok=True)

PPT_PATH = OUT_DIR / "OpenNeed-记忆稳态引擎-项目介绍-投资人版.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
PX_W = 1600
PX_H = 900

BG = (8, 12, 28)
BG2 = (15, 20, 46)
CYAN = (0, 240, 255)
BLUE = (0, 130, 255)
PURPLE = (132, 56, 255)
PINK = (255, 51, 153)
LIME = (102, 255, 204)
WHITE = (245, 248, 255)
MUTED = (163, 177, 214)


def hex_rgb(value: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*value)


def make_gradient_bg(path: Path, seed: int, accent: tuple[int, int, int], accent_2: tuple[int, int, int]) -> None:
    random.seed(seed)
    base = Image.new("RGB", (PX_W, PX_H), BG)
    px = base.load()
    for y in range(PX_H):
        t = y / max(PX_H - 1, 1)
        r = int(BG[0] * (1 - t) + BG2[0] * t)
        g = int(BG[1] * (1 - t) + BG2[1] * t)
        b = int(BG[2] * (1 - t) + BG2[2] * t)
        for x in range(PX_W):
            drift = int(10 * sin((x / PX_W) * pi * 2))
            px[x, y] = (max(0, min(255, r + drift)), max(0, min(255, g + drift)), max(0, min(255, b + drift)))

    glow = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(glow)
    for _ in range(8):
        cx = random.randint(100, PX_W - 100)
        cy = random.randint(50, PX_H - 50)
        radius = random.randint(120, 360)
        color = accent if random.random() > 0.4 else accent_2
        alpha = random.randint(40, 90)
        draw.ellipse((cx - radius, cy - radius, cx + radius, cy + radius), fill=(*color, alpha))
    glow = glow.filter(ImageFilter.GaussianBlur(55))

    overlay = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
    odraw = ImageDraw.Draw(overlay)

    for i in range(24):
        x = random.randint(0, PX_W)
        y = random.randint(0, PX_H)
        length = random.randint(180, 520)
        angle = random.uniform(-0.8, 0.8)
        x2 = x + int(length * cos(angle))
        y2 = y + int(length * sin(angle))
        color = accent if i % 2 == 0 else accent_2
        odraw.line((x, y, x2, y2), fill=(*color, 65), width=random.randint(1, 3))

    for i in range(14):
        cx = random.randint(100, PX_W - 100)
        cy = random.randint(100, PX_H - 100)
        radius = random.randint(60, 200)
        color = accent_2 if i % 2 else accent
        odraw.arc((cx - radius, cy - radius, cx + radius, cy + radius), 15, 310, fill=(*color, 120), width=3)

    overlay = overlay.filter(ImageFilter.GaussianBlur(1))

    final = Image.alpha_composite(base.convert("RGBA"), glow)
    final = Image.alpha_composite(final, overlay)

    vignette = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
    vdraw = ImageDraw.Draw(vignette)
    for i in range(16):
        inset = i * 18
        alpha = int(10 + i * 5)
        vdraw.rectangle((inset, inset, PX_W - inset, PX_H - inset), outline=(0, 0, 0, alpha), width=24)
    final = Image.alpha_composite(final, vignette)
    final.save(path)


def add_bg(slide, img_path: Path, prs: Presentation):
    slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def add_panel(slide, left, top, width, height, fill=(12, 18, 40), transparency=18, line=(80, 110, 210)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = hex_rgb(line)
    shape.line.transparency = 20
    shape.line.width = Pt(1.5)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=20, color=WHITE, bold=False, name="PingFang SC", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=20, color=WHITE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.name = "PingFang SC"
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_rgb(color)
        p.space_after = Pt(10)
    return box


def add_glow_title(slide, title, subtitle=None, accent=(0, 240, 255)):
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(11.6), Inches(0.35), "AGENT PASSPORT", 16, accent, True)
    title_box = add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.4), title, 30, WHITE, True)
    title_box.text_frame.word_wrap = True
    if subtitle:
        sub = add_textbox(slide, Inches(0.72), Inches(2.1), Inches(11.2), Inches(0.9), subtitle, 17, MUTED, False)
        sub.text_frame.word_wrap = True


def add_metric(slide, left, top, width, height, big, small, accent):
    add_panel(slide, left, top, width, height, fill=(9, 14, 34), transparency=12, line=accent)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.16), width - Inches(0.2), Inches(0.5), big, 25, accent, True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.7), width - Inches(0.25), height - Inches(0.8), small, 15, WHITE)


def build_deck():
    bg_files = []
    palettes = [
        (CYAN, PURPLE),
        (PINK, BLUE),
        (LIME, CYAN),
        (PURPLE, PINK),
    ]
    for idx in range(4):
        path = ASSET_DIR / f"bg_{idx+1}.png"
        make_gradient_bg(path, seed=42 + idx * 9, accent=palettes[idx][0], accent_2=palettes[idx][1])
        bg_files.append(path)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    def slide_with_bg(bg_idx):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, bg_files[bg_idx], prs)
        return slide

    # 1 cover
    slide = slide_with_bg(0)
    add_textbox(slide, Inches(0.75), Inches(0.68), Inches(2.2), Inches(0.35), "PROJECT INTRO", 15, CYAN, True)
    add_textbox(slide, Inches(0.75), Inches(1.2), Inches(7.8), Inches(1.8), "OpenNeed 记忆稳态引擎", 34, WHITE, True)
    add_textbox(slide, Inches(0.75), Inches(2.18), Inches(8.9), Inches(1.3), "面向 AI Agent 的身份、记忆、恢复与受控运行底座", 23, (213, 221, 255), True)
    add_textbox(slide, Inches(0.78), Inches(3.2), Inches(6.8), Inches(1.5), "让 AI 从临时聊天工具，升级为更连续、更可治理、更可恢复的运行主体。", 18, MUTED)
    add_panel(slide, Inches(8.65), Inches(1.05), Inches(3.85), Inches(4.9), fill=(12, 18, 40), transparency=8, line=CYAN)
    add_textbox(slide, Inches(8.95), Inches(1.35), Inches(3.0), Inches(0.5), "核心判断", 18, CYAN, True)
    add_textbox(slide, Inches(8.95), Inches(2.0), Inches(3.1), Inches(3.0), "未来 AI 的核心竞争力\n不只是“更聪明”\n而是“更可控、更可持续、更可治理”", 24, WHITE, True)
    add_textbox(slide, Inches(8.95), Inches(5.15), Inches(3.1), Inches(0.6), "对话不是身份\n恢复优先于重聊", 16, MUTED)

    # 2 why now
    slide = slide_with_bg(1)
    add_glow_title(slide, "为什么现在必须做 OpenNeed 记忆稳态引擎", "AI 正在从“生成内容”转向“执行任务”，身份、记忆、权限、恢复将成为新的基础设施。", PINK)
    add_metric(slide, Inches(0.75), Inches(3.0), Inches(2.8), Inches(1.5), "窗口 = 身份", "今天多数 AI 仍依附在会话里，切线程、换窗口、重启后就像换了一个人。", PINK)
    add_metric(slide, Inches(3.75), Inches(3.0), Inches(2.8), Inches(1.5), "聊天 = 记忆", "很多所谓记忆只是把更多历史塞进上下文，成本升高、质量下降、错误积累。", CYAN)
    add_metric(slide, Inches(6.75), Inches(3.0), Inches(2.8), Inches(1.5), "自动化 = 风险", "当 AI 开始调工具、改数据、碰资产，没有治理边界就无法真正进入业务流程。", LIME)
    add_metric(slide, Inches(9.75), Inches(3.0), Inches(2.8), Inches(1.5), "崩了 = 重来", "缺少恢复能力的 AI 无法承担长期工作，只能做一次性助手。", PURPLE)

    # 3 what is it
    slide = slide_with_bg(2)
    add_glow_title(slide, "OpenNeed 记忆稳态引擎是什么", "它不是普通聊天机器人，而是 AI Agent 时代的“身份层 + 运行层 + 治理层”。", LIME)
    add_panel(slide, Inches(0.75), Inches(2.2), Inches(5.75), Inches(4.55), fill=(9, 16, 34), transparency=10, line=LIME)
    add_textbox(slide, Inches(1.0), Inches(2.5), Inches(5.0), Inches(0.6), "一句话定义", 18, LIME, True)
    add_textbox(slide, Inches(1.0), Inches(3.05), Inches(4.9), Inches(1.5), "OpenNeed 记忆稳态引擎是一个让 AI Agent 拥有本地稳定身份、持续记忆、权限边界、本地可校验行为记录和恢复能力的运行底座。", 24, WHITE, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.65),
        Inches(5.0),
        Inches(1.7),
        [
            "不是“更多 Prompt”",
            "不是“更长上下文”",
            "不是“单次更聪明”",
            "而是“长期更稳定、更可恢复”",
        ],
        18,
        MUTED,
    )
    add_panel(slide, Inches(6.85), Inches(2.2), Inches(5.75), Inches(4.55), fill=(14, 18, 44), transparency=8, line=CYAN)
    add_textbox(slide, Inches(7.15), Inches(2.48), Inches(5.0), Inches(0.6), "它要替代什么旧范式", 18, CYAN, True)
    add_bullets(
        slide,
        Inches(7.15),
        Inches(3.0),
        Inches(4.9),
        Inches(3.2),
        [
            "从“对话窗口中心”升级为“身份中心”",
            "从“上下文硬撑”升级为“状态重建”",
            "从“黑箱自动化”升级为“可治理自动化”",
            "从“回答工具”升级为“长期运行底座”",
        ],
        22,
        WHITE,
    )

    # 4 problems solved
    slide = slide_with_bg(3)
    add_glow_title(slide, "它解决了什么问题", "OpenNeed 记忆稳态引擎不是做一个更会聊天的 AI，而是解决 AI 走向长期工作的底层瓶颈。", PURPLE)
    cards = [
        ("稳定身份", "关闭窗口、切换应用、重启线程之后，Agent 仍然是同一个主体。", CYAN),
        ("可恢复记忆", "忘了不是重来，而是先查本地纪要、结构化记忆、决策和恢复点。", PINK),
        ("可控执行", "低风险动作快执行，高风险动作先确认，关键动作才进入冷路径治理。", LIME),
        ("可审计责任", "为什么这么做、为什么被拦住、为什么需要人工接管，都能留下证据。", PURPLE),
    ]
    positions = [(0.8, 2.35), (6.95, 2.35), (0.8, 4.55), (6.95, 4.55)]
    for (title, body, color), (x, y) in zip(cards, positions):
        add_panel(slide, Inches(x), Inches(y), Inches(5.5), Inches(1.65), fill=(10, 16, 36), transparency=10, line=color)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.18), Inches(1.9), Inches(0.4), title, 20, color, True)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.62), Inches(4.9), Inches(0.7), body, 16, WHITE)

    # 5 architecture
    slide = slide_with_bg(0)
    add_glow_title(slide, "底层原理：五层运行架构", "身份层、记忆层、上下文重建层、验证治理层、恢复层，共同把 Agent 从会话体升级为长期运行系统。", CYAN)
    layer_titles = ["身份层", "记忆层", "上下文重建层", "验证治理层", "恢复层"]
    layer_desc = [
        "本地稳定身份、分叉关系、授权关系、本地可校验主体资料",
        "ledger / profile / episodic / working 四层记忆",
        "不拼整段聊天，按任务槽位重建当前需要的上下文",
        "模型输出必须过验证、风险分级、权限边界和必要确认",
        "checkpoint / boundary / recovery bundle / rehearsal",
    ]
    colors = [CYAN, PINK, LIME, PURPLE, BLUE]
    start_y = 2.2
    for i in range(5):
        y = start_y + i * 0.86
        add_panel(slide, Inches(0.95), Inches(y), Inches(11.3), Inches(0.65), fill=(11, 16, 37), transparency=8, line=colors[i])
        add_textbox(slide, Inches(1.22), Inches(y + 0.12), Inches(2.0), Inches(0.35), layer_titles[i], 18, colors[i], True)
        add_textbox(slide, Inches(3.05), Inches(y + 0.12), Inches(8.8), Inches(0.35), layer_desc[i], 15, WHITE)

    # 6 why local first
    slide = slide_with_bg(1)
    add_glow_title(slide, "为什么坚持 Local-First", "因为真正进入业务流程的 AI，必须先解决隐私、延迟、恢复和控制权问题。", PINK)
    add_metric(slide, Inches(0.8), Inches(2.4), Inches(2.95), Inches(1.55), "隐私", "原始记忆和高敏感状态默认不应常驻云端。", PINK)
    add_metric(slide, Inches(3.98), Inches(2.4), Inches(2.95), Inches(1.55), "低延迟", "热路径在本地闭环，才能支撑长期使用和高频协作。", CYAN)
    add_metric(slide, Inches(7.16), Inches(2.4), Inches(2.95), Inches(1.55), "可恢复", "失联、断电、重启、迁机之后仍能恢复运行。", LIME)
    add_metric(slide, Inches(10.34), Inches(2.4), Inches(2.2), Inches(1.55), "控制权", "身份和执行边界掌握在用户和组织手里。", PURPLE)
    add_panel(slide, Inches(0.8), Inches(4.45), Inches(11.75), Inches(1.85), fill=(10, 16, 34), transparency=12, line=CYAN)
    add_textbox(slide, Inches(1.05), Inches(4.72), Inches(11.0), Inches(0.4), "核心方法论", 18, CYAN, True)
    add_textbox(slide, Inches(1.05), Inches(5.18), Inches(11.1), Inches(0.8), "热路径本地化，温路径做增强，冷路径做关键动作治理。不是让所有东西都上链、都多签，而是让该快的快、该稳的稳、该审计的可审计。", 20, WHITE, True)

    # 7 use case with OpenNeed
    slide = slide_with_bg(2)
    add_glow_title(slide, "最先落地场景：招聘与人才服务", "与 OpenNeed 结合时，OpenNeed 记忆稳态引擎不只是一个后台组件，而是招聘场景里的身份连续性与受控协作底座。", LIME)
    add_bullets(
        slide,
        Inches(0.85),
        Inches(2.25),
        Inches(5.8),
        Inches(3.6),
        [
            "候选人侧：本地唯一身份、持续记忆、长期职业画像",
            "招聘顾问侧：跨岗位、跨窗口、跨线程的连续协作主体",
            "企业侧：岗位 Agent、匹配 Agent、流程 Agent 的权限边界与责任链",
            "系统侧：本地结构化处理 + 云端高质量输出的混合架构",
        ],
        22,
        WHITE,
    )
    add_panel(slide, Inches(6.95), Inches(2.25), Inches(5.35), Inches(3.75), fill=(12, 18, 40), transparency=9, line=CYAN)
    add_textbox(slide, Inches(7.25), Inches(2.55), Inches(4.8), Inches(0.5), "为什么这个场景特别适合", 20, CYAN, True)
    add_bullets(
        slide,
        Inches(7.22),
        Inches(3.1),
        Inches(4.8),
        Inches(2.5),
        [
            "招聘天然需要身份连续性",
            "多 Agent 协作天然需要授权边界",
            "候选人和岗位画像天然需要长期记忆",
            "结果输出和流程执行天然需要审计和恢复",
        ],
        18,
        MUTED,
    )

    # 8 scenarios
    slide = slide_with_bg(3)
    add_glow_title(slide, "未来使用场景", "从个人超级助手到企业级 AI 工作台，OpenNeed 记忆稳态引擎有机会成为长期运行 Agent 的基础设施。", PURPLE)
    scenario_titles = ["个人超级助手", "企业私有 AI 工作台", "高价值专业服务", "多设备 / 多 Agent 网络"]
    scenario_bodies = [
        "一个人拥有一个长期运行的本地 Agent，跨任务、跨窗口持续协助。",
        "企业内部 Agent 进入真实流程前，需要身份、权限、审计和恢复底座。",
        "法务、财务、投研、医疗辅助等高要求场景，尤其需要受控运行和审计。",
        "当单机单 Agent 做稳后，可扩展为多设备同步、团队协作和有限互通。",
    ]
    for i in range(4):
        x = 0.8 + (i % 2) * 6.15
        y = 2.2 + (i // 2) * 2.1
        color = [CYAN, PINK, LIME, BLUE][i]
        add_panel(slide, Inches(x), Inches(y), Inches(5.35), Inches(1.55), fill=(10, 16, 35), transparency=10, line=color)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.18), Inches(2.6), Inches(0.35), scenario_titles[i], 19, color, True)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.62), Inches(4.8), Inches(0.65), scenario_bodies[i], 15, WHITE)

    # 9 business model
    slide = slide_with_bg(0)
    add_glow_title(slide, "商业模式与变现路径", "先卖真实需求，再逐步形成平台层壁垒。", CYAN)
    add_metric(slide, Inches(0.8), Inches(2.4), Inches(3.8), Inches(1.7), "软件收入", "私有部署、团队版运行时、恢复能力、安全治理、运维支持。", CYAN)
    add_metric(slide, Inches(4.8), Inches(2.4), Inches(3.8), Inches(1.7), "平台收入", "Agent 身份、记忆、恢复与治理中枢，服务多个上层业务 Agent。", PINK)
    add_metric(slide, Inches(8.8), Inches(2.4), Inches(3.8), Inches(1.7), "生态收入", "插件、工具接入、企业集成、审计与协作网络。", LIME)
    add_panel(slide, Inches(0.85), Inches(4.75), Inches(11.65), Inches(1.45), fill=(12, 18, 40), transparency=10, line=PURPLE)
    add_textbox(slide, Inches(1.1), Inches(5.0), Inches(11.0), Inches(0.7), "先不卖“全网协议叙事”，先卖“可落地的 Agent Runtime”。一旦真实业务在这层沉淀，后续互通能力才有意义。", 21, WHITE, True)

    # 10 roadmap
    slide = slide_with_bg(1)
    add_glow_title(slide, "路线图", "从单机单 Agent 出发，逐步走向团队协作、企业部署和多节点网络。", PINK)
    roadmap = [
        ("Phase 1", "单机单 Agent、本地优先、可恢复、可审计", CYAN),
        ("Phase 2", "企业内多 Agent 协作与权限治理", PINK),
        ("Phase 3", "多设备同步、备份、迁移与恢复体系", LIME),
        ("Phase 4", "跨产品、跨业务的 Agent 身份互通与受控协作", PURPLE),
    ]
    for i, (phase, desc, color) in enumerate(roadmap):
        y = 2.3 + i * 1.0
        add_panel(slide, Inches(0.9), Inches(y), Inches(11.5), Inches(0.68), fill=(10, 16, 35), transparency=10, line=color)
        add_textbox(slide, Inches(1.15), Inches(y + 0.12), Inches(1.35), Inches(0.35), phase, 18, color, True)
        add_textbox(slide, Inches(2.55), Inches(y + 0.12), Inches(9.2), Inches(0.35), desc, 17, WHITE)

    # 11 moat
    slide = slide_with_bg(2)
    add_glow_title(slide, "潜在壁垒", "真正稀缺的不是“再做一个 Agent UI”，而是形成可审计 Agent 基础设施的系统能力。", LIME)
    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.25),
        Inches(11.2),
        Inches(3.8),
        [
            "产品壁垒：身份、记忆、恢复、治理、审计不是单点功能，而是系统闭环",
            "数据壁垒：一旦真实工作流沉淀在 Passport 中，迁移成本会快速上升",
            "安全壁垒：可校验、可恢复、可追责能力天然更适合企业付费",
            "生态壁垒：未来多个业务 Agent 共享同一身份与治理底座时，将形成平台效应",
            "战略壁垒：占位“长期运行 Agent”这一新对象，而不是继续卷单次模型输出",
        ],
        21,
        WHITE,
    )

    # 12 closing
    slide = slide_with_bg(3)
    add_textbox(slide, Inches(0.75), Inches(0.72), Inches(2.2), Inches(0.35), "FINAL MESSAGE", 15, PINK, True)
    add_textbox(slide, Inches(0.75), Inches(1.25), Inches(11.2), Inches(1.2), "下一代 AI 的关键，不是“更会聊天”。", 31, WHITE, True)
    add_textbox(slide, Inches(0.75), Inches(2.15), Inches(11.6), Inches(1.2), "而是第一次拥有“可持续恢复、可持续治理、可持续运行”的能力。", 31, WHITE, True)
    add_panel(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.45), fill=(11, 18, 40), transparency=8, line=CYAN)
    add_textbox(slide, Inches(1.05), Inches(4.3), Inches(11.0), Inches(0.7), "OpenNeed 记忆稳态引擎想做的，就是这个时代的 Agent 身份层、记忆层和受控运行底座。", 24, CYAN, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.85), Inches(6.55), Inches(4.8), Inches(0.3), "OpenNeed 记忆稳态引擎 · Investor Deck", 13, MUTED, True)

    prs.save(PPT_PATH)
    print(PPT_PATH)


if __name__ == "__main__":
    build_deck()
